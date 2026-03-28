from __future__ import annotations

from functools import cached_property
from typing import Any

PRESENTATIONML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


class PptxConverter:
    def to_markdown(self, path: str) -> str:
        payload = self.to_json(path)
        parts: list[str] = []

        if payload["title"]:
            parts.append(f"# {payload['title']}")
            parts.append("")

        for slide in payload["slides"]:
            parts.append(f"## Slide {slide['slide_number']}")
            if slide["title"]:
                parts.append(f"### {slide['title']}")

            for block in slide["blocks"]:
                block_type = block["type"]
                if block_type == "text":
                    parts.append(block["text"])
                elif block_type == "bullets":
                    for item in block["items"]:
                        indent = "  " * max(item["level"], 0)
                        parts.append(f"{indent}- {item['text']}")
                elif block_type == "table":
                    parts.extend(self._table_to_markdown(block["rows"]))

            if slide["notes"]:
                parts.append("### Notes")
                parts.append(slide["notes"])

            parts.append("")

        markdown = "\n".join(parts).strip()
        return f"{markdown}\n" if markdown else ""

    def to_json(self, path: str) -> dict[str, Any]:
        presentation = self._presentation_class(path)
        slides = [
            self._slide_to_payload(slide, index + 1)
            for index, slide in enumerate(presentation.slides)
        ]

        return {
            "title": self._resolve_title(presentation, slides),
            "slides": slides,
            "metadata": {
                "subject": self._normalize_text(presentation.core_properties.subject),
                "author": self._normalize_text(presentation.core_properties.author),
                "keywords": self._normalize_text(presentation.core_properties.keywords),
            },
        }

    def _slide_to_payload(self, slide, slide_number: int) -> dict[str, Any]:
        title_shape = getattr(slide.shapes, "title", None)
        title = self._normalize_text(getattr(title_shape, "text", None))
        blocks: list[dict[str, Any]] = []

        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue

            if getattr(shape, "has_table", False):
                payload = self._table_block(shape.table)
            elif getattr(shape, "has_text_frame", False):
                payload = self._text_block(shape)
            else:
                payload = None

            if payload is not None:
                blocks.append(payload)

        notes = self._extract_notes(slide)

        return {
            "slide_number": slide_number,
            "title": title,
            "blocks": blocks,
            "notes": notes,
        }

    def _text_block(self, shape) -> dict[str, Any] | None:
        paragraphs: list[dict[str, Any]] = []
        bullet_like = False

        for paragraph in getattr(shape.text_frame, "paragraphs", []):
            text = self._normalize_text(getattr(paragraph, "text", None))
            if not text:
                continue
            level = int(getattr(paragraph, "level", 0) or 0)
            bullet_like = bullet_like or self._is_bullet_paragraph(paragraph)
            paragraphs.append(
                {
                    "text": text,
                    "level": level,
                }
            )

        if not paragraphs:
            return None

        if len(paragraphs) == 1:
            return {
                "type": "text",
                "text": paragraphs[0]["text"],
            }

        if not bullet_like:
            return {
                "type": "text",
                "text": "\n\n".join(paragraph["text"] for paragraph in paragraphs),
            }

        return {
            "type": "bullets",
            "items": paragraphs,
        }

    def _is_bullet_paragraph(self, paragraph) -> bool:
        level = int(getattr(paragraph, "level", 0) or 0)
        if level > 0:
            return True

        properties = getattr(paragraph, "_pPr", None)
        if properties is None:
            return False

        return any(
            properties.find(path, namespaces=PRESENTATIONML_NS) is not None
            for path in ("a:buChar", "a:buAutoNum", "a:buBlip")
        )

    def _table_block(self, table) -> dict[str, Any] | None:
        rows: list[list[str]] = []

        for row in getattr(table, "rows", []):
            values = [self._normalize_text(cell.text) for cell in getattr(row, "cells", [])]
            if any(values):
                rows.append(values)

        if not rows:
            return None

        return {
            "type": "table",
            "rows": rows,
        }

    def _extract_notes(self, slide) -> str:
        if not getattr(slide, "has_notes_slide", False):
            return ""

        notes_slide = slide.notes_slide
        text_frame = getattr(notes_slide, "notes_text_frame", None)
        return self._normalize_text(getattr(text_frame, "text", None))

    def _resolve_title(self, presentation, slides: list[dict[str, Any]]) -> str:
        core_title = self._normalize_text(presentation.core_properties.title)
        if core_title:
            return core_title

        for slide in slides:
            if slide["title"]:
                return slide["title"]

        return ""

    def _table_to_markdown(self, rows: list[list[str]]) -> list[str]:
        if not rows:
            return []
        width = max(len(row) for row in rows)
        padded_rows = [row + [""] * max(0, width - len(row)) for row in rows]
        header = padded_rows[0]
        body = padded_rows[1:]

        lines = [
            "| " + " | ".join(self._escape_md(value) for value in header) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]

        for row in body:
            lines.append("| " + " | ".join(self._escape_md(value) for value in row) + " |")

        return lines

    def _normalize_text(self, value: Any) -> str:
        if value is None:
            return ""
        return str(value).strip()

    def _escape_md(self, text: str) -> str:
        return text.replace("|", "\\|").replace("\n", " ").replace("\r", " ").strip()

    @cached_property
    def _presentation_class(self):
        from pptx import Presentation

        return Presentation
