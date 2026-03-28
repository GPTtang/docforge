#!/usr/bin/env python3
"""Generate richer Office fixtures for DocForge regression testing."""

from __future__ import annotations

import base64
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Font
from openpyxl.worksheet.hyperlink import Hyperlink
from pptx import Presentation
from pptx.util import Inches as PptxInches


REPO_ROOT = Path(__file__).resolve().parents[1]
FIXTURES = REPO_ROOT / "test-documents"

PNG_1X1_RED = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WnR7VQAAAAASUVORK5CYII="
)


def ensure_dirs() -> None:
    for relative in [
        "word/docx",
        "powerpoint/pptx",
        "excel/xlsx",
        "assets",
    ]:
        (FIXTURES / relative).mkdir(parents=True, exist_ok=True)


def write_support_image() -> Path:
    image_path = FIXTURES / "assets" / "tiny-red.png"
    image_path.write_bytes(PNG_1X1_RED)
    return image_path


def build_docx_main(image_path: Path) -> None:
    path = FIXTURES / "word" / "docx" / "unit_test_formatting.docx"
    doc = Document()
    doc.core_properties.title = "DocForge Word Regression Fixture"
    doc.core_properties.subject = "Word structure regression"
    doc.core_properties.author = "Codex"
    doc.core_properties.keywords = "docx, regression, headings, tables"

    doc.add_heading("DocForge Word Regression Fixture", level=1)
    intro = doc.add_paragraph()
    intro.add_run("Purpose: ").bold = True
    intro.add_run("validate headings, paragraphs, tables, lists, and inline formatting.")

    doc.add_heading("Structured Content", level=2)
    para = doc.add_paragraph()
    para.add_run("This paragraph mixes ")
    para.add_run("bold").bold = True
    para.add_run(", ")
    para.add_run("italic").italic = True
    para.add_run(", and ")
    para.add_run("underlined").underline = True
    para.add_run(" text.")

    doc.add_paragraph("Primary bullet", style="List Bullet")
    doc.add_paragraph("Nested bullet", style="List Bullet 2")
    doc.add_paragraph("Ordered item", style="List Number")

    doc.add_heading("Quarterly Snapshot", level=2)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ["Metric", "Q1", "Q2"]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    rows = [
        ("Revenue", "120", "135"),
        ("Cost", "80", "82"),
        ("Status", "Stable", "Improving"),
    ]
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = value

    doc.add_heading("Embedded Image", level=2)
    image_para = doc.add_paragraph("Minimal image fixture:")
    image_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    doc.add_picture(str(image_path), width=Inches(0.35))

    doc.add_heading("Bilingual Segment", level=2)
    doc.add_paragraph("English heading with 中文段落 and 日本語メモ to catch Unicode regressions.")

    doc.save(path)


def build_docx_table() -> None:
    path = FIXTURES / "word" / "docx" / "tablecell.docx"
    doc = Document()
    doc.core_properties.title = "DocForge Table Cell Fixture"
    doc.add_heading("Merged Table Fixture", level=1)

    table = doc.add_table(rows=4, cols=4)
    table.style = "Table Grid"
    table.cell(0, 0).merge(table.cell(0, 1)).text = "Region"
    table.cell(0, 2).merge(table.cell(0, 3)).text = "Result"
    table.cell(1, 0).text = "APAC"
    table.cell(1, 1).text = "Tokyo"
    table.cell(1, 2).text = "Pass"
    table.cell(1, 3).text = "92%"
    table.cell(2, 0).text = "EMEA"
    table.cell(2, 1).text = "Berlin"
    table.cell(2, 2).text = "Warn"
    table.cell(2, 3).text = "74%"
    table.cell(3, 0).merge(table.cell(3, 3)).text = "Footer: merged row"

    doc.save(path)


def build_pptx_main() -> None:
    path = FIXTURES / "powerpoint" / "pptx" / "powerpoint_sample.pptx"
    prs = Presentation()
    prs.core_properties.title = "DocForge PowerPoint Regression Fixture"
    prs.core_properties.subject = "Slides with bullets and tables"
    prs.core_properties.author = "Codex"
    prs.core_properties.keywords = "pptx, regression, bullets, table"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Launch Overview"
    body = slide.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = "Key deliverables"
    p.level = 0
    p = body.add_paragraph()
    p.text = "API compatibility maintained"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Regression suite expanded"
    p.level = 1
    slide.notes_slide.notes_text_frame.text = "Presenter note: emphasize local-first architecture."

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Scorecard"
    table_shape = slide2.shapes.add_table(3, 3, PptxInches(1), PptxInches(1.8), PptxInches(6), PptxInches(1.5))
    table = table_shape.table
    values = [
        ["Metric", "Before", "After"],
        ["PDF", "Docling/Marker", "OpenDataLoader"],
        ["Office", "Docling", "Specialized parsers"],
    ]
    for r_idx, row in enumerate(values):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = value
    slide2.notes_slide.notes_text_frame.text = "Table slide for markdown table rendering."

    prs.save(path)


def build_pptx_image(image_path: Path) -> None:
    path = FIXTURES / "powerpoint" / "pptx" / "powerpoint_with_image.pptx"
    prs = Presentation()
    prs.core_properties.title = "DocForge Image Slide Fixture"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Image and Summary"
    textbox = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.6), PptxInches(4.5), PptxInches(1.5))
    text_frame = textbox.text_frame
    text_frame.text = "Compact image fixture"
    p = text_frame.add_paragraph()
    p.text = "Used to verify text + image coexistence."
    p.level = 0
    slide.shapes.add_picture(str(image_path), PptxInches(5.6), PptxInches(1.8), width=PptxInches(1.0))
    slide.notes_slide.notes_text_frame.text = "Image should not break text extraction."

    prs.save(path)


def build_xlsx_main() -> None:
    path = FIXTURES / "excel" / "xlsx" / "xlsx_01.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "Metric"
    ws["B1"] = "Value"
    ws["A2"] = "Revenue"
    ws["B2"] = 135
    ws["A3"] = "Cost"
    ws["B3"] = 82
    ws["A4"] = "Margin"
    ws["B4"] = "=B2-B3"
    ws["A5"] = "Reference"
    ws["B5"] = "DocForge"
    ws["B5"].hyperlink = Hyperlink(ref="B5", target="https://example.com/docforge")
    ws["B5"].font = Font(color="0000EE", underline="single")
    ws["A6"] = "Commented"
    ws["B6"] = "Check benchmark"
    ws["B6"].comment = Comment("Generated fixture comment", "Codex")
    ws["A7"] = "Merged"
    ws.merge_cells("B7:C7")
    ws["B7"] = "Merged range"

    detail = wb.create_sheet("Detail")
    detail.append(["Quarter", "North", "South"])
    detail.append(["Q1", 44, 41])
    detail.append(["Q2", 46, 43])
    detail.append(["Q3", 48, 45])

    wb.save(path)


def build_xlsx_table() -> None:
    path = FIXTURES / "excel" / "xlsx" / "xlsx_05_table_with_title.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws["A1"] = "DocForge Sheet Fixture"
    ws["A1"].font = Font(bold=True, size=14)
    ws.merge_cells("A1:D1")
    ws.append(["Item", "Owner", "Score", "Status"])
    ws.append(["Parser", "Platform", 91, "green"])
    ws.append(["DOCX", "Word", 88, "amber"])
    ws.append(["PPTX", "Slides", 86, "amber"])
    ws.append(["XLSX", "Sheets", "=AVERAGE(C3:C4)", "derived"])
    ws["D6"].comment = Comment("Formula row", "Codex")

    wb.save(path)


def main() -> None:
    ensure_dirs()
    image_path = write_support_image()
    build_docx_main(image_path)
    build_docx_table()
    build_pptx_main()
    build_pptx_image(image_path)
    build_xlsx_main()
    build_xlsx_table()
    print("Generated Office fixtures under test-documents/")


if __name__ == "__main__":
    main()
